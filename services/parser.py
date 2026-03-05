from __future__ import annotations

import re
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from models.slides import Deck, Slide


# ─── Helpers ──────────────────────────────────────────────────────────────────

def _clean(text: str) -> str:
    """Strip excess whitespace and non-printable characters."""
    return re.sub(r"\s+", " ", text).strip()


def _extract_texts(shape) -> list[str]:
    """Return all non-empty paragraph texts from a shape's text frame."""
    if not shape.has_text_frame:
        return []
    return [_clean(p.text) for p in shape.text_frame.paragraphs if _clean(p.text)]


def _is_title_shape(shape) -> bool:
    try:
        from pptx.util import Emu
        from pptx.enum.shapes import PP_PLACEHOLDER
        ph = shape.placeholder_format
        return ph is not None and ph.idx in (0, 1)  # 0=title, 1=center title
    except Exception:
        return False


def _is_body_shape(shape) -> bool:
    try:
        ph = shape.placeholder_format
        return ph is not None and ph.idx >= 2
    except Exception:
        return False


# ─── Parser ───────────────────────────────────────────────────────────────────

def parse_pptx(file_bytes: bytes) -> Deck:
    """
    Parse a .pptx binary into a Deck.

    Strategy per slide:
    - Title: first placeholder with idx 0 or 1
    - Subtitle: second title-like placeholder OR first body paragraph if short
    - Bullets: remaining text paragraphs from body placeholders
    - Notes: notes slide text
    """
    import io
    prs = Presentation(io.BytesIO(file_bytes))
    slides: list[Slide] = []

    for idx, pptx_slide in enumerate(prs.slides):
        title = ""
        subtitle = ""
        bullets: list[str] = []
        notes = ""

        # ── Collect shape texts ────────────────────────────────────────────
        title_texts: list[str] = []
        body_texts: list[str] = []

        for shape in pptx_slide.shapes:
            if not shape.has_text_frame:
                continue

            texts = _extract_texts(shape)
            if not texts:
                continue

            if _is_title_shape(shape):
                title_texts.extend(texts)
            elif _is_body_shape(shape):
                body_texts.extend(texts)
            else:
                # Non-placeholder text boxes — treat as body
                body_texts.extend(texts)

        # ── Assign title / subtitle ────────────────────────────────────────
        if title_texts:
            title = title_texts[0]
            if len(title_texts) > 1:
                subtitle = title_texts[1]

        # If still no subtitle, check if first body line is short enough
        if not subtitle and body_texts and len(body_texts[0]) < 80:
            subtitle = body_texts.pop(0)

        bullets = body_texts

        # Fallback: if no title was found at all, use first bullet
        if not title and bullets:
            title = bullets.pop(0)

        # ── Notes ─────────────────────────────────────────────────────────
        try:
            notes_slide = pptx_slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes = _clean(notes_tf.text) if notes_tf else ""
        except Exception:
            notes = ""

        slides.append(
            Slide(
                id=idx,
                title=title or f"Slide {idx + 1}",
                subtitle=subtitle,
                bullets=bullets,
                notes=notes,
            )
        )

    return Deck(slides=slides)